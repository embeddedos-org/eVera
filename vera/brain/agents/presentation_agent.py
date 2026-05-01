"""Presentation Agent -- PowerPoint slides, pitch decks, PDF export."""

from __future__ import annotations

import logging
from typing import Any

from vera.brain.agents.base import BaseAgent, Tool
from vera.providers.models import ModelTier

logger = logging.getLogger(__name__)


class CreateSlidesTool(Tool):
    def __init__(self):
        super().__init__(
            name="create_slides",
            description="Create PowerPoint presentation",
            parameters={
                "title": {"type": "str", "description": "Title"},
                "slides": {"type": "str", "description": "JSON array [{title,content,notes}]"},
                "output": {"type": "str", "description": "Output .pptx path"},
            },
        )

    async def execute(self, **kw: Any) -> dict[str, Any]:
        try:
            import json

            from pptx import Presentation
            from pptx.util import Inches

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = kw.get("title", "Presentation")
            slides = json.loads(kw.get("slides", "[]")) if isinstance(kw.get("slides"), str) else kw.get("slides", [])
            for s in slides:
                sl = prs.slides.add_slide(prs.slide_layouts[1])
                sl.shapes.title.text = s.get("title", "")
                if len(sl.placeholders) > 1:
                    sl.placeholders[1].text = s.get("content", "")
                if s.get("notes"):
                    sl.notes_slide.notes_text_frame.text = s["notes"]
            out = kw.get("output", "presentation.pptx")
            prs.save(out)
            return {"status": "success", "output": out, "slides": len(slides) + 1}
        except ImportError:
            return {"status": "error", "message": "pip install python-pptx"}
        except Exception as e:
            return {"status": "error", "message": str(e)}


class SlideTemplateTool(Tool):
    def __init__(self):
        super().__init__(
            name="slide_templates",
            description="Get presentation templates",
            parameters={
                "type": {
                    "type": "str",
                    "description": "pitch_deck|quarterly_review|project_proposal|training|sales|keynote",
                }
            },
        )

    async def execute(self, **kw: Any) -> dict[str, Any]:
        templates = {
            "pitch_deck": ["Title", "Problem", "Solution", "Market Size", "Business Model", "Traction", "Team", "Ask"],
            "quarterly_review": [
                "Title",
                "Summary",
                "Metrics",
                "Revenue",
                "Achievements",
                "Challenges",
                "Goals",
                "Q&A",
            ],
            "project_proposal": [
                "Title",
                "Background",
                "Objectives",
                "Scope",
                "Timeline",
                "Budget",
                "Risks",
                "Conclusion",
            ],
            "training": ["Title", "Agenda", "Objectives", "Topic 1", "Topic 2", "Exercise", "Summary", "Q&A"],
            "sales": ["Title", "Challenge", "Solution", "Features", "Case Studies", "Pricing", "ROI", "Next Steps"],
            "keynote": ["Title", "Vision", "Story", "Demo", "Impact", "Future", "Call to Action", "Thank You"],
        }
        return {
            "status": "success",
            "type": kw.get("type", "pitch_deck"),
            "structure": templates.get(kw.get("type", "pitch_deck"), templates["pitch_deck"]),
        }


class ExportPDFTool(Tool):
    def __init__(self):
        super().__init__(
            name="export_slides_pdf",
            description="Export slides to PDF",
            parameters={
                "input": {"type": "str", "description": "Input .pptx"},
                "output": {"type": "str", "description": "Output .pdf"},
            },
        )

    async def execute(self, **kw: Any) -> dict[str, Any]:
        import subprocess

        try:
            subprocess.run(["libreoffice", "--headless", "--convert-to", "pdf", kw["input"]], capture_output=True)
            return {"status": "success", "output": kw.get("output", kw["input"].replace(".pptx", ".pdf"))}
        except Exception as e:
            return {"status": "error", "message": str(e)}


class PresentationAgent(BaseAgent):
    name = "presentation"
    description = "Create PowerPoint slides, pitch decks, templates, PDF export"
    tier = ModelTier.SPECIALIST
    system_prompt = "You are eVera's Presentation Agent. Create PowerPoint presentations, use templates, export to PDF."
    offline_responses = {
        "slide": "\U0001f4ca Creating slides!",
        "presentation": "\U0001f3ac Building!",
        "deck": "\U0001f4ca Making deck!",
        "pitch": "\U0001f680 Pitch deck!",
    }

    def _setup_tools(self):
        self._tools = [CreateSlidesTool(), SlideTemplateTool(), ExportPDFTool()]
